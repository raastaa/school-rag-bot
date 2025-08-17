# ingest/ingest_generic.py
from __future__ import annotations
import os
import re
import uuid
import argparse
from typing import List, Dict, Any, Tuple

from dotenv import load_dotenv

load_dotenv()

from pypdf import PdfReader

try:
    from docx import Document as Docx

    DOCX_AVAILABLE = True
except Exception:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation

    PPTX_AVAILABLE = True
except Exception:
    PPTX_AVAILABLE = False

try:
    import openpyxl

    XLSX_AVAILABLE = True
except Exception:
    XLSX_AVAILABLE = False

import zipfile
from xml.etree import ElementTree as ET
from bs4 import BeautifulSoup
import trafilatura

# локальные утилиты
from .pdf_ingest import (
    read_pdf_pages as _read_pdf_pages,
)  # можно использовать готовый из pdf_ingest
from text_utils import (
    normalize_pdf_text,
    split_text_hard,
    split_into_chunks,
    count_tokens,
    EMB_MAX,
)
from gigachat_client import GigaChatEmbedder, detect_dim
from store_qdrant import ensure_collection, upsert_chunks
from db_local import has_chunk_hash, save_chunk_hash
import hashlib

SOURCE_GROUP_DEFAULT = "zabedu"

# -------------------- readers --------------------


def read_pdf(file_path: str) -> List[Tuple[int, str]]:
    """
    Возвращает [(page_number, text)], начиная с 1.
    Пытается аккуратно нормализовать текст (склейка переносов с дефисами и т.п.).
    """
    try:
        reader = PdfReader(file_path)
        pages: List[Tuple[int, str]] = []
        for i, page in enumerate(reader.pages, start=1):
            try:
                raw = page.extract_text() or ""
            except Exception:
                raw = ""
            pages.append((i, normalize_pdf_text(raw)))
        return pages
    except Exception:
        # fallback — через готовую функцию (если в вашем проекте она надёжнее)
        return _read_pdf_pages(file_path)


def read_docx(file_path: str) -> str:
    """
    Читает DOCX. Если нет python-docx — парсит как zip -> word/document.xml (fallback).
    """
    # 1) Нормальный путь через python-docx
    if DOCX_AVAILABLE:
        try:
            d = Docx(file_path)
            paras = [p.text for p in d.paragraphs]
            text = "\n".join(paras).strip()
            if text:
                return text
        except Exception:
            pass

    # 2) Fallback: DOCX как ZIP
    try:
        with zipfile.ZipFile(file_path) as z:
            with z.open("word/document.xml") as f:
                xml_bytes = f.read()
        # Простое извлечение текста из XML (без форматирования)
        root = ET.fromstring(xml_bytes)
        # В DOCX текст в <w:t>. Нужны пространства имён.
        ns = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
        texts = [el.text for el in root.findall(".//w:t", ns) if el.text]
        return "\n".join(texts).strip()
    except Exception:
        return ""


def read_pptx(file_path: str) -> str:
    if not PPTX_AVAILABLE:
        return ""
    try:
        prs = Presentation(file_path)
        texts: List[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
        return "\n".join(texts).strip()
    except Exception:
        return ""


def read_xlsx(file_path: str) -> str:
    if not XLSX_AVAILABLE:
        return ""
    try:
        wb = openpyxl.load_workbook(file_path, read_only=True)
        texts: List[str] = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                row_txt = " ".join(str(c) for c in row if c is not None).strip()
                if row_txt:
                    texts.append(row_txt)
        return "\n".join(texts).strip()
    except Exception:
        return ""


def read_html(file_path: str) -> str:
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            html = f.read()
        txt = (
            trafilatura.extract(html, include_formatting=False, include_links=False)
            or ""
        )
        if not txt:
            soup = BeautifulSoup(html, "html.parser")
            for s in soup(["script", "style", "noscript"]):
                s.decompose()
            txt = re.sub(r"\s+", " ", soup.get_text(" ").strip())
        return txt
    except Exception:
        return ""


def read_txt(file_path: str) -> str:
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    except Exception:
        return ""


# -------------------- core ingest --------------------


async def ingest_file(
    file_path: str,
    title: str | None = None,
    source_group: str = SOURCE_GROUP_DEFAULT,
    doc_tag: str | None = None,
    heading_path: str | None = None,
) -> Dict[str, Any]:
    """
    Индексирует один файл в Qdrant. Поддержка: .pdf .docx .html/.htm .txt
    Возвращает {'document', 'chunks', 'pages'(pdf), 'file'}
    """
    ext = os.path.splitext(file_path)[1].lower()
    abs_path = os.path.abspath(file_path)

    # 1) размерность модели эмбеддингов и ensure коллекции
    dim = await detect_dim()
    ensure_collection(dim)

    # 2) извлекаем текст/страницы и формируем чанки
    chunks: List[Tuple[str, int | None, int | None]] = []  # (text, page_from, page_to)
    pages = 0

    if ext == ".pdf":
        pages_list = read_pdf(file_path)  # [(p, text)]
        pages = len(pages_list)
        # делим с учётом EMB_MAX + нормализованных страниц (page_from/page_to корректны)
        chunks = split_into_chunks(pages_list)
    elif ext == ".docx":
        t = read_docx(file_path)
        if not t:
            return {
                "document": title or os.path.basename(file_path),
                "chunks": 0,
                "pages": 0,
                "file": abs_path,
            }
        parts = split_text_hard(t, EMB_MAX)
        chunks = [(p, None, None) for p in parts]
    elif ext in (".html", ".htm"):
        t = read_html(file_path)
        if not t:
            return {
                "document": title or os.path.basename(file_path),
                "chunks": 0,
                "pages": 0,
                "file": abs_path,
            }
        parts = split_text_hard(t, EMB_MAX)
        chunks = [(p, None, None) for p in parts]
    elif ext == ".txt":
        t = read_txt(file_path)
        if not t:
            return {
                "document": title or os.path.basename(file_path),
                "chunks": 0,
                "pages": 0,
                "file": abs_path,
            }
        parts = split_text_hard(t, EMB_MAX)
        chunks = [(p, None, None) for p in parts]
    elif ext == ".pptx":
        t = read_pptx(file_path)
        if not t:
            return {
                "document": title or os.path.basename(file_path),
                "chunks": 0,
                "pages": 0,
                "file": abs_path,
            }
        parts = split_text_hard(t, EMB_MAX)
        chunks = [(p, None, None) for p in parts]
    elif ext in (".xlsx", ".xlsm"):
        t = read_xlsx(file_path)
        if not t:
            return {
                "document": title or os.path.basename(file_path),
                "chunks": 0,
                "pages": 0,
                "file": abs_path,
            }
        parts = split_text_hard(t, EMB_MAX)
        chunks = [(p, None, None) for p in parts]
    else:
        # неизвестный тип — пропускаем
        return {
            "document": title or os.path.basename(file_path),
            "chunks": 0,
            "pages": 0,
            "file": abs_path,
        }

    # 3) формируем payloads/texts для записи в векторное хранилище
    texts: List[str] = []
    payloads: List[Dict[str, Any]] = []
    seq = 0
    doc_title = title or os.path.basename(file_path)
    token_total = 0

    for text_c, p_from, p_to in chunks:
        h = hashlib.sha256(text_c.strip().encode("utf-8")).hexdigest()
        if has_chunk_hash(h):
            continue
        save_chunk_hash(h, abs_path, seq)
        tok = count_tokens(text_c)
        payloads.append(
            {
                "id": str(uuid.uuid4()),
                "seq": seq,
                "source": os.path.basename(file_path),
                "source_group": source_group,
                "title": doc_title,
                "page_from": p_from,
                "page_to": p_to,
                "path": abs_path,
                "type": ext.lstrip("."),
                "token_count": tok,
                "text": text_c[:1000],
                "doc_tag": doc_tag,
                "heading_path": heading_path,
            }
        )
        texts.append(text_c)
        token_total += tok
        seq += 1

    # 4) эмбеддинги + upsert в Qdrant
    embedder = GigaChatEmbedder()
    total = 0
    BATCH = int(os.getenv("EMBEDDING_BATCH", "64"))
    for i in range(0, len(texts), BATCH):
        batch_vecs = await embedder.embed(texts[i : i + BATCH])
        upsert_chunks(batch_vecs, payloads[i : i + BATCH])
        total += len(batch_vecs)

    return {
        "document": doc_title,
        "chunks": total,
        "pages": pages,
        "file": abs_path,
        "source_group": source_group,
        "doc_tag": doc_tag,
        "heading_path": heading_path,
        "token_total": token_total,
    }


async def ingest_path(
    root: str,
    source_group: str = SOURCE_GROUP_DEFAULT,
    exts: Tuple[str, ...] = (
        ".pdf",
        ".docx",
        ".html",
        ".htm",
        ".txt",
        ".pptx",
        ".xlsx",
        ".xlsm",
    ),
    doc_tag: str | None = None,
) -> Dict[str, int]:
    """
    Рекурсивно обходит директорию и индексирует все поддерживаемые файлы.
    Возвращает счётчики {'files': N, 'chunks': M}.
    """
    total_files = 0
    total_chunks = 0
    for dirpath, _, filenames in os.walk(root):
        rel = os.path.relpath(dirpath, root)
        heading_path = None if rel in ("", ".") else " > ".join(rel.split(os.sep))
        for name in filenames:
            if os.path.splitext(name)[1].lower() in exts:
                total_files += 1
                fp = os.path.join(dirpath, name)
                try:
                    res = await ingest_file(
                        fp,
                        title=name,
                        source_group=source_group,
                        doc_tag=doc_tag,
                        heading_path=heading_path,
                    )
                    total_chunks += res.get("chunks", 0)
                    print(
                        f"[OK] {name}: chunks={res.get('chunks')} pages={res.get('pages')}"
                    )
                except Exception as e:
                    print(f"[FAIL] {name}: {e}")
    return {"files": total_files, "chunks": total_chunks}


# -------------------- CLI --------------------

if __name__ == "__main__":
    import asyncio

    ap = argparse.ArgumentParser()
    ap.add_argument(
        "--path", required=True, help="Папка с файлами (распакованный архив)"
    )
    ap.add_argument(
        "--group",
        default=SOURCE_GROUP_DEFAULT,
        help="source_group для Qdrant (по умолчанию zabedu)",
    )
    ap.add_argument("--doc-tag", default=None, help="doc_tag payload value")
    args = ap.parse_args()

    print("Indexing:", args.path, "group:", args.group)
    res = asyncio.run(ingest_path(args.path, source_group=args.group, doc_tag=args.doc_tag))
    print("Done:", res)
